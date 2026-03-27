#!/usr/bin/env python3
"""Test script to create a simple PPT presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.core_properties.title = "Hello PPT"

# Add title slide (layout index 0)
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)

# Set title
title = slide.shapes.title
title.text = "Hello PPT"

# Save
output_path = "/Users/arctic/Desktop/Hello_PPT.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
