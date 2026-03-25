#!/usr/bin/env python3
"""
PPTX Maker - Python MCP Server

A Model Context Protocol (MCP) server for creating and editing PowerPoint presentations.
Based on the mcp-builder template.

Usage:
    python server.py

Requirements:
    pip install mcp python-pptx Pillow

This server provides tools for:
- Creating new PowerPoint presentations
- Adding slides with various layouts
- Adding text, bullet points, images, and shapes
- Saving presentations to files
- Generating slide previews
- Opening presentations in desktop apps
"""

import json
import os
import sys
import subprocess
import tempfile
from pathlib import Path
from typing import Any, List

try:
    from mcp.server import Server
    from mcp.types import Tool, TextContent
    from mcp.server.stdio import stdio_server
except ImportError:
    print("Error: MCP library not installed. Run: pip install mcp", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("Error: Pillow not installed. Run: pip install Pillow", file=sys.stderr)
    sys.exit(1)

# Initialize the MCP server
server = Server("pptx-maker")

# Temp file for state persistence between tool calls
TEMP_DIR = "/tmp/_pptx_maker"
TEMP_PRESENTATION = f"{TEMP_DIR}/current.pptx"
TEMP_PREVIEWS_DIR = f"{TEMP_DIR}/previews"

# Ensure temp directories exist
os.makedirs(TEMP_DIR, exist_ok=True)
os.makedirs(TEMP_PREVIEWS_DIR, exist_ok=True)


def save_to_temp(prs):
    """Save presentation to temp file for persistence."""
    prs.save(TEMP_PRESENTATION)


def load_from_temp():
    """Load presentation from temp file if it exists."""
    if os.path.exists(TEMP_PRESENTATION):
        try:
            return Presentation(TEMP_PRESENTATION)
        except Exception:
            pass
    return None


def result_text(data: dict) -> list[TextContent]:
    """Helper to create success response."""
    return [TextContent(type="text", text=json.dumps({"success": True, **data}))]


def error_text(msg: str) -> list[TextContent]:
    """Helper to create error response."""
    return [TextContent(type="text", text=json.dumps({"error": msg}))]


def render_slide_to_image(slide, width=960, height=540) -> Image.Image:
    """Render a slide to a PIL Image (simplified rendering)."""
    # Create white background image
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)

    # Try to get a font
    try:
        font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
        font_body = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
    except Exception:
        font_title = ImageFont.load_default()
        font_body = ImageFont.load_default()

    # Draw title
    if hasattr(slide, 'shapes') and slide.shapes.title:
        title_text = slide.shapes.title.text if slide.shapes.title.text else ""
        if title_text:
            draw.text((40, 40), title_text[:100], fill='black', font=font_title)

    # Draw a placeholder for content (simplified)
    draw.rectangle([(40, 100), (width - 40, height - 40)], outline='lightgray', width=1)
    draw.text((50, height // 2 - 20), "[Slide Content]", fill='gray', font=font_body)

    # Draw slide number indicator
    slide_num = f"Slide Preview ({width}x{height})"
    draw.text((width - 120, height - 30), slide_num, fill='lightgray', font=font_body)

    return img


@server.list_tools()
async def list_tools() -> list[Tool]:
    """
    Return the list of available tools.
    """
    return [
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation. This clears any existing presentation.",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Presentation title"},
                    "author": {"type": "string", "description": "Author name"}
                },
                "required": ["title"]
            }
        ),
        Tool(
            name="add_slide",
            description="Add a new slide to the presentation. Creates a slide with the specified layout.",
            inputSchema={
                "type": "object",
                "properties": {
                    "layout": {
                        "type": "string",
                        "description": "Slide layout type",
                        "enum": ["title", "title_content", "two_content", "blank"]
                    },
                    "title": {"type": "string", "description": "Slide title (optional)"}
                }
            }
        ),
        Tool(
            name="add_text",
            description="Add text to the current slide.",
            inputSchema={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Text content to add"},
                    "position": {
                        "type": "string",
                        "description": "Position: 'title' or 'content'",
                        "enum": ["title", "content"]
                    },
                    "font_size": {"type": "number", "description": "Font size in points"},
                    "bold": {"type": "boolean", "description": "Bold text"}
                },
                "required": ["text"]
            }
        ),
        Tool(
            name="add_bullet_points",
            description="Add bullet points to the current slide.",
            inputSchema={
                "type": "object",
                "properties": {
                    "items": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of bullet point items"
                    }
                },
                "required": ["items"]
            }
        ),
        Tool(
            name="add_image",
            description="Add an image to the current slide.",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to the image file"},
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"}
                },
                "required": ["path"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Save the current presentation to a file path. Returns the saved path.",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Output file path (e.g., /path/to/presentation.pptx)"}
                },
                "required": ["path"]
            }
        ),
        Tool(
            name="get_info",
            description="Get information about the current or a specific presentation.",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PPTX file (optional, uses temp if not provided)"}
                }
            }
        ),
        Tool(
            name="preview_slides",
            description="Generate PNG preview images of slides. Returns base64-encoded image data.",
            inputSchema={
                "type": "object",
                "properties": {
                    "slide_numbers": {
                        "type": "array",
                        "items": {"type": "number"},
                        "description": "List of slide numbers to preview (1-indexed). If empty, previews all slides."
                    },
                    "width": {
                        "type": "number",
                        "description": "Preview width in pixels",
                        "default": 480
                    },
                    "height": {
                        "type": "number",
                        "description": "Preview height in pixels",
                        "default": 270
                    }
                }
            }
        ),
        Tool(
            name="open_presentation",
            description="Open the presentation in the default desktop application (PowerPoint, Keynote, etc.).",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path to PPTX file (optional, uses temp if not provided)"}
                }
            }
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict[str, Any]) -> list[TextContent]:
    """
    Handle tool execution requests.
    """
    try:
        if name == "create_presentation":
            # Clean up old previews
            for f in os.listdir(TEMP_PREVIEWS_DIR):
                try:
                    os.remove(os.path.join(TEMP_PREVIEWS_DIR, f))
                except Exception:
                    pass

            # Remove existing temp file
            if os.path.exists(TEMP_PRESENTATION):
                os.remove(TEMP_PRESENTATION)

            prs = Presentation()
            prs.core_properties.title = arguments.get("title", "Untitled")
            prs.core_properties.author = arguments.get("author", "")
            save_to_temp(prs)

            return result_text({
                "message": f"Created presentation: {arguments['title']}",
                "note": "Use add_slide to add slides"
            })

        elif name == "add_slide":
            prs = load_from_temp()
            if not prs:
                return error_text("No presentation found. Use create_presentation first.")

            layout_map = {
                "title": 0,
                "title_content": 1,
                "two_content": 2,
                "blank": 6,
            }
            layout_name = arguments.get("layout", "title_content")
            layout_idx = layout_map.get(layout_name, 1)

            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

            slide_title = arguments.get("title", "")
            if slide_title and hasattr(slide, 'shapes') and slide.shapes.title:
                slide.shapes.title.text = slide_title

            save_to_temp(prs)

            return result_text({
                "slide_number": len(prs.slides),
                "layout": layout_name,
                "title": slide_title
            })

        elif name == "add_text":
            prs = load_from_temp()
            if not prs or not prs.slides:
                return error_text("No slide available. Use add_slide first.")

            slide = prs.slides[-1]
            text = arguments["text"]
            position = arguments.get("position", "content")
            font_size = arguments.get("font_size", 18)
            bold = arguments.get("bold", False)

            if position == "title" and hasattr(slide, 'shapes') and slide.shapes.title:
                slide.shapes.title.text = text
            else:
                textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                tf = textbox.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold

            save_to_temp(prs)
            return result_text({"text": text})

        elif name == "add_bullet_points":
            prs = load_from_temp()
            if not prs or not prs.slides:
                return error_text("No slide available. Use add_slide first.")

            slide = prs.slides[-1]
            items = arguments["items"]

            # Find or create content placeholder
            content_shape = None
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape != getattr(slide.shapes, 'title', None):
                    content_shape = shape
                    break

            if not content_shape:
                content_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))

            tf = content_shape.text_frame
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {item}"

            save_to_temp(prs)
            return result_text({"items_added": len(items)})

        elif name == "add_image":
            prs = load_from_temp()
            if not prs or not prs.slides:
                return error_text("No slide available. Use add_slide first.")

            image_path = arguments["path"]
            if not os.path.exists(image_path):
                return error_text(f"Image not found: {image_path}")

            slide = prs.slides[-1]
            left = Inches(arguments.get("left", 1))
            top = Inches(arguments.get("top", 2))
            width = Inches(arguments.get("width", 6))

            slide.shapes.add_picture(image_path, left, top, width=width)
            save_to_temp(prs)

            return result_text({"path": image_path})

        elif name == "save_presentation":
            prs = load_from_temp()
            if not prs:
                return error_text("No presentation to save. Use create_presentation first.")

            path = arguments["path"]
            Path(path).parent.mkdir(parents=True, exist_ok=True)

            prs.save(path)
            save_to_temp(prs)

            return result_text({
                "path": path,
                "slides": len(prs.slides),
                "open_url": f"/api/pptx/open?path={path}"
            })

        elif name == "get_info":
            path = arguments.get("path") or TEMP_PRESENTATION

            if os.path.exists(path):
                prs = Presentation(path)
                info = {
                    "path": path,
                    "slide_count": len(prs.slides),
                    "title": prs.core_properties.title or "Untitled",
                    "author": prs.core_properties.author or "Unknown",
                }
            else:
                info = {"path": path, "exists": False}

            return [TextContent(type="text", text=json.dumps(info, indent=2))]

        elif name == "preview_slides":
            path = arguments.get("path") or TEMP_PRESENTATION

            if not os.path.exists(path):
                return error_text("No presentation found. Use create_presentation first.")

            prs = Presentation(path)
            slide_numbers = arguments.get("slide_numbers", [])
            width = arguments.get("width", 480)
            height = arguments.get("height", 270)

            if not slide_numbers:
                slide_numbers = list(range(1, len(prs.slides) + 1))

            previews = []
            for idx in slide_numbers:
                if 1 <= idx <= len(prs.slides):
                    slide = prs.slides[idx - 1]
                    img = render_slide_to_image(slide, width, height)

                    # Save to temp and return path
                    preview_path = f"{TEMP_PREVIEWS_DIR}/slide_{idx}.png"
                    img.save(preview_path, "PNG")

                    # Return base64 data
                    import base64
                    with open(preview_path, "rb") as f:
                        img_data = base64.b64encode(f.read()).decode()

                    previews.append({
                        "slide_number": idx,
                        "width": width,
                        "height": height,
                        "data": f"data:image/png;base64,{img_data}"
                    })

            return result_text({
                "previews": previews,
                "total_slides": len(prs.slides),
                "previewed_slides": len(previews)
            })

        elif name == "open_presentation":
            path = arguments.get("path") or TEMP_PRESENTATION

            if not os.path.exists(path):
                return error_text(f"File not found: {path}")

            try:
                # macOS: use 'open' command
                subprocess.run(["open", path], check=True)
                return result_text({
                    "message": f"Opened: {path}",
                    "path": path
                })
            except subprocess.CalledProcessError as e:
                return error_text(f"Failed to open file: {str(e)}")
            except FileNotFoundError:
                return error_text("'open' command not found. This feature works on macOS.")

        else:
            return error_text(f"Unknown tool: {name}")

    except Exception as e:
        return error_text(str(e))


async def main():
    """Run the MCP server using stdio transport."""
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options()
        )


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
